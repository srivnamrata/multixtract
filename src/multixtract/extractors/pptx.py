"""Native PowerPoint (.pptx) extractor.

Ported from the ZF PPT + GPT-4o Vision pipeline. Extracts, with python-pptx:
  * presentation metadata (slide count + dimensions)
  * per-slide title, text (with GroupShape recursion + SmartArt), tables,
    and hyperlinks (single pass over slide shapes)
  * images from ``ppt/media/`` mapped to their slide via slide relationships

Vector images (EMF/WMF/SVG, plus EMF-in-.bin) are converted to PNG in one
headless LibreOffice batch call; JPEG XR (.wdp) via imagecodecs when available.
Images are filtered through the shared :class:`ImageFilterPipeline` and returned
as ``prepared_images`` for downstream vision analysis.

Requires the ``[pptx]`` extra (python-pptx). EMF/WMF/SVG need system LibreOffice;
.wdp needs optional ``imagecodecs``. Both degrade gracefully.
"""
from __future__ import annotations

import io
import logging
import os
import xml.etree.ElementTree as ET
import zipfile
from collections import defaultdict
from typing import Any, Dict, List, Optional, Set, Tuple

from ..filters import ImageFilterPipeline
from ._image_utils import (
    IMAGE_EXTS,
    VECTOR_EXTS,
    WDP_EXTS,
    batch_convert_vectors_to_png,
    decode_wdp_to_png,
    ensure_rgb_png,
)

log = logging.getLogger("multixtract.extractors.pptx")

_A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
_PKG_REL_NS = "http://schemas.openxmlformats.org/package/2006/relationships"


def _iter_all_shapes(shapes, mso):
    """Recursively yield shapes, descending into GroupShapes."""
    for shape in shapes:
        if shape.shape_type == mso.GROUP:
            yield from _iter_all_shapes(shape.shapes, mso)
        else:
            yield shape


def _extract_smartart_text(shape) -> Optional[str]:
    """Extract SmartArt text from drawingml <a:t> nodes."""
    try:
        root = ET.fromstring(shape.element.xml)
        texts = [(t.text or "").strip() for t in root.iter(f"{{{_A_NS}}}t")]
        texts = [t for t in texts if t]
        if texts:
            return " | ".join(texts)
    except Exception:
        pass
    return None


def _extract_slide_content(slide, mso) -> Tuple[str, str, List[List[List[str]]], List[str]]:
    """Single-pass slide extraction: (text, title, tables, hyperlinks)."""
    texts: List[str] = []
    title = ""
    tables: List[List[List[str]]] = []
    hyperlinks: List[str] = []

    for shape in _iter_all_shapes(slide.shapes, mso):
        shape_type = shape.shape_type
        if shape_type == mso.EMBEDDED_OLE_OBJECT:
            continue

        if shape.has_text_frame:
            parts = []
            for para in shape.text_frame.paragraphs:
                pt = para.text.strip()
                if pt:
                    parts.append(pt)
                for run in para.runs:
                    try:
                        if run.hyperlink and run.hyperlink.address:
                            hyperlinks.append(run.hyperlink.address)
                    except Exception:
                        pass
            shape_text = "\n".join(parts)
            if not shape_text:
                continue
            try:
                if shape.is_placeholder and shape.placeholder_format.idx == 0:
                    title = shape_text
            except Exception:
                pass
            texts.append(shape_text)

        elif shape.has_table:
            table_data = [[cell.text.strip() for cell in row.cells] for row in shape.table.rows]
            if table_data:
                tables.append(table_data)

        elif shape_type == mso.PICTURE:
            pass  # images handled separately
        else:
            smartart_text = _extract_smartart_text(shape)
            if smartart_text:
                texts.append(f"[SmartArt] {smartart_text}")

    return "\n".join(texts), title, tables, hyperlinks


def _build_slide_media_map(zf: zipfile.ZipFile, n_slides: int) -> Dict[int, List[str]]:
    """Map slide number -> [media paths] from each slide's relationships."""
    slide_media: Dict[int, List[str]] = defaultdict(list)
    namelist = set(zf.namelist())
    for slide_num in range(1, n_slides + 1):
        rels_path = f"ppt/slides/_rels/slide{slide_num}.xml.rels"
        if rels_path not in namelist:
            continue
        try:
            root = ET.fromstring(zf.read(rels_path).decode("utf-8"))
            for rel in root.findall(f"{{{_PKG_REL_NS}}}Relationship"):
                target = rel.get("Target", "")
                if "../media/" in target:
                    slide_media[slide_num].append(f"ppt/media/{target.replace('../media/', '')}")
        except Exception as exc:
            log.debug("rels parse failed for slide %d: %s", slide_num, exc)
    return slide_media


def _looks_like_emf_bin(raw: bytes) -> bool:
    """Detect EMF metafiles stored with a .bin extension."""
    return len(raw) >= 44 and raw[:4] == b"\x01\x00\x00\x00" and raw[40:44] == b" EMF"


class PptxExtractor:
    """DocumentExtractor for ``.pptx`` (native python-pptx extraction)."""

    extensions: Tuple[str, ...] = (".pptx",)

    def __init__(self, vector_timeout: int = 120) -> None:
        self.vector_timeout = vector_timeout

    def extract(
        self,
        path: str,
        image_filter: Optional[ImageFilterPipeline] = None,
    ) -> Tuple[Dict[str, Any], List[Dict[str, Any]]]:
        try:
            from pptx import Presentation
            from pptx.enum.shapes import MSO_SHAPE_TYPE as mso
        except ImportError as e:
            raise ImportError(
                "PowerPoint support requires python-pptx: pip install 'multixtract[pptx]'"
            ) from e
        from PIL import Image

        if image_filter is None:
            image_filter = ImageFilterPipeline()
        image_filter.reset()

        base_name = os.path.splitext(os.path.basename(path))[0]
        empty: Dict[str, Any] = {"_base_name": base_name, "metadata": {}, "pgs": []}
        try:
            prs = Presentation(path)
            n_slides = len(prs.slides)

            document: Dict[str, Any] = {
                "metadata": {
                    "slide_count":  n_slides,
                    "slide_width":  prs.slide_width,
                    "slide_height": prs.slide_height,
                },
                "_base_name": base_name,
                "pgs": [],
            }

            prepared_images: List[Dict[str, Any]] = []
            try:
                zf = zipfile.ZipFile(path, "r")
            except Exception as exc:
                log.warning("could not open %s as ZIP for image extraction: %s", base_name, exc)
                # Still return slide text even if the archive can't be opened.
                for idx, slide in enumerate(prs.slides):
                    txt, title, tables, links = _extract_slide_content(slide, mso)
                    document["pgs"].append({
                        "pg_num": idx + 1, "kind": "slide", "title": title,
                        "txt": txt, "tables": tables,
                        "hyperlinks": list(dict.fromkeys(links)), "imgs": [],
                    })
                return document, prepared_images

            try:
                slide_media_map = _build_slide_media_map(zf, n_slides)
                seen_media: Set[str] = set()

                # Pre-scan: collect vector / WDP images that need conversion.
                vector_items, wdp_items = [], []
                for slide_num in range(1, n_slides + 1):
                    for media_path in slide_media_map.get(slide_num, []):
                        if media_path in seen_media:
                            continue
                        ext = os.path.splitext(media_path)[1].lower()
                        try:
                            raw = zf.read(media_path)
                        except KeyError:
                            continue
                        if ext in VECTOR_EXTS:
                            vector_items.append((media_path, raw))
                            seen_media.add(media_path)
                        elif ext in WDP_EXTS:
                            wdp_items.append((media_path, raw))
                            seen_media.add(media_path)
                        elif ext == ".bin" and _looks_like_emf_bin(raw):
                            vector_items.append((media_path, raw))
                            seen_media.add(media_path)
                converted = batch_convert_vectors_to_png(vector_items, self.vector_timeout)
                converted.update(decode_wdp_to_png(wdp_items))

                # Track media paths already sent to prepare_image (deduplicates
                # converted vectors that appear on multiple slides — Bug 8).
                processed_media: Set[str] = set()

                # Process slides: content + per-slide images.
                for slide_idx, slide in enumerate(prs.slides):
                    slide_num = slide_idx + 1
                    txt, title, tables, links = _extract_slide_content(slide, mso)
                    document["pgs"].append({
                        "pg_num":     slide_num,
                        "kind":       "slide",
                        "title":      title,
                        "txt":        txt,
                        "tables":     tables,
                        "hyperlinks": list(dict.fromkeys(links)),
                        "imgs":       [],
                    })

                    img_idx = 0
                    for media_path in slide_media_map.get(slide_num, []):
                        ext = os.path.splitext(media_path)[1].lower()
                        if ext not in IMAGE_EXTS:
                            continue

                        # Deduplicate across all slides (covers both converted vectors
                        # and rasters that appear in multiple slide relationships).
                        if media_path in processed_media:
                            image_filter.note_duplicate()
                            continue

                        if media_path in converted:
                            image_bytes, ext_out = converted[media_path], "png"
                        elif ext in VECTOR_EXTS or ext in WDP_EXTS or ext == ".bin":
                            continue  # conversion failed or non-image .bin
                        else:
                            if media_path in seen_media:
                                image_filter.note_duplicate()
                                continue  # raster already handled on an earlier slide
                            try:
                                image_bytes = zf.read(media_path)
                            except KeyError:
                                continue
                            seen_media.add(media_path)
                            ext_out = ext.lstrip(".")
                            ext_out = {"tif": "tiff", "jpg": "jpeg"}.get(ext_out, ext_out)
                            if ext_out == "png" and not image_bytes[:4].startswith(b"\x89PNG"):
                                fixed = ensure_rgb_png(image_bytes)
                                if fixed is None:
                                    continue
                                image_bytes = fixed

                        processed_media.add(media_path)

                        try:
                            with Image.open(io.BytesIO(image_bytes)) as image:
                                width, height = image.size
                        except Exception as exc:
                            log.debug("image decode failed for %s on slide %d of %s: %s",
                                      media_path, slide_num, base_name, exc)
                            continue

                        prepared = image_filter.prepare_image(
                            image_bytes=image_bytes,
                            ext=ext_out,
                            width=width,
                            height=height,
                            image_id=f"page_{slide_num}_img_{img_idx}",
                            page_number=slide_num,
                            img_idx=img_idx,
                        )
                        if prepared is not None:
                            prepared_images.append(prepared)
                            img_idx += 1  # only advance for images that pass filtering
            finally:
                zf.close()

            return document, prepared_images
        except ImportError:
            raise
        except Exception:
            log.warning("PptxExtractor failed for %s", path, exc_info=True)
            return empty, []
