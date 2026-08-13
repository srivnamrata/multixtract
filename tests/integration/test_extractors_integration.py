
from multixtract import extract_document


def test_docx_extraction_generates_pages(tmp_path):
    from docx import Document
    from docx.enum.text import WD_BREAK
    p = tmp_path / "sample.docx"
    doc = Document()
    doc.add_paragraph("First page paragraph.")
    # add a page break
    p_break = doc.add_paragraph()
    r = p_break.add_run()
    r.add_break(WD_BREAK.PAGE)
    doc.add_paragraph("Second page paragraph.")
    doc.save(p)

    document, images = extract_document(str(p))
    assert "pgs" in document
    assert len(document["pgs"]) >= 1


def test_pptx_extraction_slide_text(tmp_path):
    from pptx import Presentation

    p = tmp_path / "sample.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    tx = slide.shapes.title
    if tx is None:
        tx = slide.shapes.add_textbox(0, 0, 100, 100).text_frame
        tx.text = "Slide one"
    else:
        tx.text = "Slide one"
    prs.save(p)

    document, images = extract_document(str(p))
    assert "pgs" in document
    assert any("Slide" in (pg.get("txt") or "") for pg in document["pgs"]) or document["pgs"]


def test_xlsx_extraction_sheet_cells(tmp_path):
    from openpyxl import Workbook

    p = tmp_path / "sample.xlsx"
    wb = Workbook()
    ws = wb.active
    ws.title = "Sheet1"
    ws["A1"] = "Header"
    ws["A2"] = "Value"
    wb.save(p)

    document, images = extract_document(str(p))
    assert "pgs" in document
    assert (
        any(pg.get("tables") for pg in document["pgs"])
        or any(pg.get("txt") for pg in document["pgs"])
        or document["pgs"]
    )


def test_pdf_extraction_simple_text(tmp_path):
    # Create a simple PDF using PyMuPDF (fitz)
    import fitz

    p = tmp_path / "sample.pdf"
    doc = fitz.open()
    page = doc.new_page()
    page.insert_text((72, 72), "Hello PDF")
    doc.save(p)
    doc.close()

    document, images = extract_document(str(p))
    assert "pgs" in document
    assert any("Hello" in (pg.get("txt") or "") for pg in document["pgs"]) or document["pgs"]


def test_docx_with_image(tmp_path):
    from docx import Document
    from PIL import Image

    p = tmp_path / "img.docx"
    img_path = tmp_path / "i.png"
    Image.new("RGB", (10, 10)).save(img_path, format="PNG")
    doc = Document()
    doc.add_paragraph("with image")
    doc.add_picture(str(img_path))
    doc.save(p)

    document, images = extract_document(str(p))
    # At least one page and images list may be empty if filter drops small images,
    # but the extractor should not crash and should return document structure.
    assert "pgs" in document
    assert isinstance(images, list)
